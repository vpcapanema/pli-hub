#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Gera a apresentacao do portfolio PLI Hub em HTML e PPTX a partir de dados.json.

Uso:  python gerar_apresentacao.py

Saidas (na mesma pasta):
  - pli-hub-apresentacao.html  (deck navegavel; imagens embutidas em base64)
  - pli-hub-apresentacao.pptx  (PowerPoint 16:9)

O conteudo vive APENAS em dados.json e em imagens/. Para mudar um texto ou
trocar uma captura, edite o JSON (ou o arquivo em imagens/) e rode este
script de novo; HTML e PPTX nunca divergem.
"""

import base64
import json
import pathlib

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

AQUI = pathlib.Path(__file__).parent
DADOS = json.loads((AQUI / "dados.json").read_text(encoding="utf-8"))

# ---------------------------------------------------------------- paleta

GROUND = RGBColor(0x06, 0x1A, 0x26)
SURFACE = RGBColor(0x0C, 0x29, 0x39)
SURFACE2 = RGBColor(0x10, 0x32, 0x46)
RULE = RGBColor(0x1A, 0x46, 0x5F)
TEXT = RGBColor(0xE8, 0xF1, 0xF6)
MUTED = RGBColor(0x8F, 0xB0, 0xC2)
DIM = RGBColor(0x5D, 0x82, 0x96)
BLUE = RGBColor(0x3A, 0xA6, 0xD9)
GREEN = RGBColor(0x3E, 0xC2, 0x6E)
AMBER = RGBColor(0xE0, 0xA3, 0x3A)
DARKINK = RGBColor(0x0A, 0x2A, 0x16)

DISPLAY = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGEM = Inches(0.72)
LARG = W - 2 * MARGEM

# geometria da captura no slide de aplicacao
SHOT_X = Inches(5.95)
SHOT_W = W - MARGEM - SHOT_X          # ~6.66"
SHOT_H = Inches(3.05)
SHOT_Y = Inches(1.28)
ASPECTO = float(SHOT_W) / float(SHOT_H)


def hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def cor_cat(cat):
    return hex_rgb(DADOS["categorias"][cat]["cor"])


def recorte(caminho, aspecto):
    """Recorta a captura no aspecto do slide (a partir do topo) e devolve o arquivo."""
    origem = AQUI / caminho
    destino = AQUI / "imagens" / "_pptx" / origem.name
    destino.parent.mkdir(parents=True, exist_ok=True)
    if destino.exists() and destino.stat().st_mtime >= origem.stat().st_mtime:
        return str(destino)
    im = Image.open(origem).convert("RGB")
    alvo = int(im.width / aspecto)
    im = im.crop((0, 0, im.width, min(alvo, im.height)))
    im.save(destino, "JPEG", quality=84, optimize=True)
    return str(destino)


# ---------------------------------------------------------------- helpers


def novo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GROUND
    return slide


def faixa(slide, cor):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(4))
    s.fill.solid()
    s.fill.fore_color.rgb = cor
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def caixa(slide, x, y, w, h, preenche=None, borda=None, larg_borda=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if preenche is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = preenche
    if borda is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = borda
        s.line.width = larg_borda
    s.shadow.inherit = False
    return s


def texto(slide, x, y, w, h, conteudo, tam=14, cor=TEXT, fonte=BODY, negrito=False,
          espaco=1.0, alinha=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, linha in enumerate(conteudo if isinstance(conteudo, list) else [conteudo]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinha
        p.line_spacing = espaco
        if i:
            p.space_before = Pt(6)
        r = p.add_run()
        r.text = linha
        r.font.size = Pt(tam)
        r.font.color.rgb = cor
        r.font.name = fonte
        r.font.bold = negrito
    return tb


def rodape(slide, esq, dir_, n, total):
    y = H - Inches(0.55)
    linha = caixa(slide, MARGEM, y - Inches(0.15), LARG, Pt(0.75), preenche=RULE)
    linha.line.fill.background()
    texto(slide, MARGEM, y, Inches(7), Inches(0.28), esq, tam=8.5, cor=DIM, fonte=MONO)
    texto(slide, W - MARGEM - Inches(5), y, Inches(5), Inches(0.28),
          f"{dir_}     {n:02d}/{total}", tam=8.5, cor=DIM, fonte=MONO, alinha=PP_ALIGN.RIGHT)


def eyebrow(slide, x, y, txt, cor=BLUE, w=Inches(6)):
    return texto(slide, x, y, w, Inches(0.24), txt.upper(), tam=9, cor=cor, fonte=MONO)


def pilula(slide, x, y, txt, fundo, cor_txt):
    largura = Inches(0.085) * len(txt) + Inches(0.34)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, largura, Inches(0.26))
    s.fill.solid()
    s.fill.fore_color.rgb = fundo
    s.line.fill.background()
    s.shadow.inherit = False
    s.adjustments[0] = 0.15
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt.upper()
    r.font.size = Pt(8.5)
    r.font.name = MONO
    r.font.color.rgb = cor_txt
    return s


def chips(slide, x, y, w, itens):
    cx, cy = x, y
    for chip in itens:
        largura = Inches(0.062) * len(chip) + Inches(0.26)
        if cx + largura > x + w:
            cx, cy = x, cy + Inches(0.34)
        c = caixa(slide, cx, cy, largura, Inches(0.26), preenche=None, borda=RULE)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = chip
        r.font.size = Pt(8)
        r.font.name = MONO
        r.font.color.rgb = MUTED
        cx += largura + Inches(0.09)
    return cy + Inches(0.26)


def lista_destaques(slide, x, y, w, itens, cor, tam=10.5, chars=74):
    for d in itens:
        linhas = max(1, -(-len(d) // chars))
        alt = Inches(0.2) * linhas
        t = caixa(slide, x, y + Inches(0.09), Inches(0.14), Pt(1.2), preenche=cor)
        t.line.fill.background()
        texto(slide, x + Inches(0.28), y, w - Inches(0.28), alt, d,
              tam=tam, cor=RGBColor(0xD3, 0xE3, 0xEC), espaco=1.18)
        y += alt + Inches(0.11)
    return y


# ---------------------------------------------------------------- slides


def slide_capa(prs, n, total):
    s = novo_slide(prs)
    faixa(s, BLUE)
    caixa(s, W - Inches(5.4), 0, Inches(5.4), H, preenche=RGBColor(0x08, 0x22, 0x30))
    caixa(s, W - Inches(5.4), 0, Pt(1), H, preenche=RULE)

    eyebrow(s, MARGEM, Inches(0.95), f"{DADOS['org']} · {DADOS['data']}")
    texto(s, MARGEM, Inches(1.5), Inches(7.4), Inches(2.2), "PLI Hub",
          tam=88, negrito=True, fonte=DISPLAY, espaco=0.85)
    texto(s, MARGEM, Inches(3.75), Inches(6.2), Inches(1.4), DADOS["subtitulo"],
          tam=21, cor=MUTED, espaco=1.2)

    y = Inches(5.05)
    for rotulo, valor in [
        ("Aplicações", f"{len(DADOS['apps'])} sistemas em 3 categorias"),
        ("Central de acesso", DADOS["hub_url"]),
        ("Responsável técnico", DADOS["autor"]),
    ]:
        texto(s, MARGEM, y, Inches(6.5), Inches(0.24), rotulo.upper(), tam=9, cor=DIM, fonte=MONO)
        texto(s, MARGEM, y + Inches(0.23), Inches(6.5), Inches(0.3), valor, tam=13)
        y += Inches(0.62)

    x = W - Inches(5.4) + Inches(0.7)
    cw = Inches(4.0)
    yy = Inches(1.35)
    texto(s, x, yy, cw, Inches(0.26), "NESTE DOCUMENTO", tam=9, cor=DIM, fonte=MONO)
    yy += Inches(0.45)
    for cid, cat in DADOS["categorias"].items():
        caixa(s, x, yy, cw, Pt(2.5), preenche=hex_rgb(cat["cor"]))
        texto(s, x, yy + Inches(0.12), cw, Inches(0.28), cat["nome"],
              tam=13.5, negrito=True, fonte=DISPLAY, cor=hex_rgb(cat["cor"]))
        yy += Inches(0.42)
        for a in [a for a in DADOS["apps"] if a["categoria"] == cid]:
            texto(s, x, yy, cw, Inches(0.24), a["nome"], tam=11, cor=MUTED)
            yy += Inches(0.26)
        yy += Inches(0.2)
    return s


def slide_panorama(prs, n, total):
    s = novo_slide(prs)
    faixa(s, BLUE)
    apps = DADOS["apps"]
    no_ar = sum(1 for a in apps if a["status"] == "no ar")
    com_shot = sum(1 for a in apps if a.get("imagem"))

    eyebrow(s, MARGEM, Inches(0.5), "Panorama")
    texto(s, MARGEM, Inches(0.88), Inches(9), Inches(0.7), "O portfólio em um relance",
          tam=34, negrito=True, fonte=DISPLAY)

    x = MARGEM
    for valor, rotulo, cor in [
        (str(len(apps)), "aplicações", TEXT), ("3", "categorias", TEXT),
        (str(no_ar), "no ar hoje", GREEN), (str(len(apps) - no_ar), "suspensas", AMBER),
        (str(com_shot), "com tela registrada", TEXT),
    ]:
        texto(s, x, Inches(1.75), Inches(2.4), Inches(0.7), valor, tam=38, negrito=True, fonte=DISPLAY, cor=cor)
        texto(s, x, Inches(2.35), Inches(2.4), Inches(0.3), rotulo.upper(), tam=8.5, cor=DIM, fonte=MONO)
        x += Inches(2.42)
    caixa(s, MARGEM, Inches(2.8), LARG, Pt(0.75), preenche=RULE)

    col_w = (LARG - Inches(0.7)) / 3
    x = MARGEM
    for cid, cat in DADOS["categorias"].items():
        lista = [a for a in apps if a["categoria"] == cid]
        cor = hex_rgb(cat["cor"])
        caixa(s, x, Inches(3.15), col_w, Pt(3), preenche=cor)
        texto(s, x, Inches(3.35), col_w, Inches(0.6), cat["nome"],
              tam=16, negrito=True, fonte=DISPLAY, cor=cor, espaco=1.05)
        texto(s, x, Inches(3.95), col_w, Inches(0.25), f"{len(lista)} aplicações",
              tam=9, cor=DIM, fonte=MONO)
        y = Inches(4.35)
        for k, a in enumerate(lista, 1):
            texto(s, x, y, Inches(0.3), Inches(0.26), f"{k:02d}", tam=9, cor=DIM, fonte=MONO)
            texto(s, x + Inches(0.42), y - Inches(0.02), col_w - Inches(0.8), Inches(0.26), a["nome"], tam=12.5)
            ponto = caixa(s, x + col_w - Inches(0.22), y + Inches(0.05), Inches(0.1), Inches(0.1),
                          preenche=GREEN if a["status"] == "no ar" else AMBER)
            ponto.line.fill.background()
            caixa(s, x, y + Inches(0.3), col_w - Inches(0.1), Pt(0.75), preenche=RULE)
            y += Inches(0.44)
        x += col_w + Inches(0.35)

    rodape(s, DADOS["hub_url"], "Panorama", n, total)
    return s


def slide_app(prs, app, i, total_apps, n, total):
    s = novo_slide(prs)
    cor = cor_cat(app["categoria"])
    cat = DADOS["categorias"][app["categoria"]]["nome"]
    faixa(s, cor)

    eyebrow(s, MARGEM, Inches(0.5), f"Aplicação {i + 1:02d} / {total_apps}", cor=cor, w=Inches(2.2))
    texto(s, MARGEM + Inches(2.3), Inches(0.5), Inches(5), Inches(0.24), cat.upper(),
          tam=9, cor=DIM, fonte=MONO)
    if app["status"] == "no ar":
        pilula(s, W - MARGEM - Inches(0.85), Inches(0.45), "no ar", GREEN, DARKINK)
    else:
        pilula(s, W - MARGEM - Inches(1.1), Inches(0.45), "suspensa", AMBER, RGBColor(0x2A, 0x14, 0x08))

    # ------- coluna esquerda: identidade
    esq_w = Inches(4.9)
    texto(s, MARGEM, Inches(1.22), esq_w, Inches(0.8), app["nome"], tam=33, negrito=True, fonte=DISPLAY, espaco=0.95)
    texto(s, MARGEM, Inches(1.95), esq_w, Inches(0.5), app["extenso"], tam=11.5, cor=MUTED, espaco=1.15)
    texto(s, MARGEM, Inches(2.72), esq_w, Inches(1.7), app["resumo"], tam=12.5, espaco=1.3)

    caixa(s, MARGEM, Inches(4.62), Pt(3), Inches(0.4), preenche=cor)
    texto(s, MARGEM + Inches(0.15), Inches(4.67), esq_w, Inches(0.36), app["papel"],
          tam=13.5, negrito=True, fonte=DISPLAY, cor=cor)

    caixa(s, MARGEM, Inches(5.3), esq_w, Pt(0.75), preenche=RULE)
    chips(s, MARGEM, Inches(5.5), esq_w, app["stack"])

    # ------- coluna direita: captura + entregas
    if app.get("imagem"):
        arq = recorte(app["imagem"], ASPECTO)
        barra = caixa(s, SHOT_X, SHOT_Y - Inches(0.26), SHOT_W, Inches(0.26), preenche=SURFACE2, borda=RULE)
        tf = barra.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = tf.paragraphs[0].add_run()
        r.text = "   " + app["url"]
        r.font.size = Pt(8)
        r.font.name = MONO
        r.font.color.rgb = DIM
        for k in range(3):
            d = caixa(s, SHOT_X + Inches(0.1) + k * Inches(0.14), SHOT_Y - Inches(0.2),
                      Inches(0.08), Inches(0.08), preenche=RULE)
            d.line.fill.background()
        s.shapes.add_picture(arq, SHOT_X, SHOT_Y, width=SHOT_W, height=SHOT_H)
        moldura = caixa(s, SHOT_X, SHOT_Y, SHOT_W, SHOT_H, preenche=None, borda=RULE)
        moldura.line.width = Pt(0.75)
        texto(s, SHOT_X, SHOT_Y + SHOT_H + Inches(0.08), SHOT_W, Inches(0.24),
              app["imagem_legenda"], tam=8, cor=DIM, fonte=MONO)
        y_entrega = SHOT_Y + SHOT_H + Inches(0.5)
    else:
        alerta = caixa(s, SHOT_X, SHOT_Y, SHOT_W, Inches(1.15), preenche=None, borda=RULE)
        alerta.line.dash_style = 4  # tracejado
        texto(s, SHOT_X + Inches(0.3), SHOT_Y + Inches(0.22), SHOT_W - Inches(0.6), Inches(0.28),
              "Sem captura disponível", tam=13, negrito=True, fonte=DISPLAY, cor=AMBER)
        texto(s, SHOT_X + Inches(0.3), SHOT_Y + Inches(0.55), SHOT_W - Inches(0.6), Inches(0.5),
              "A aplicação está suspensa no Render e não aparece nas capturas do relatório de "
              "atividades. A tela entra assim que o serviço for reativado.",
              tam=10.5, cor=DIM, espaco=1.2)
        y_entrega = SHOT_Y + Inches(1.6)

    texto(s, SHOT_X, y_entrega, SHOT_W, Inches(0.24), "O QUE ELA ENTREGA", tam=8.5, cor=DIM, fonte=MONO)
    lista_destaques(s, SHOT_X, y_entrega + Inches(0.34), SHOT_W, app["destaques"], cor)

    rodape(s, app["url"], cat, n, total)
    return s


def slide_modulos(prs, n, total):
    s = novo_slide(prs)
    faixa(s, BLUE)
    eyebrow(s, MARGEM, Inches(0.5), "SIGMA-PLI · por dentro")
    texto(s, MARGEM, Inches(0.88), Inches(10), Inches(0.7), "A plataforma central, módulo a módulo",
          tam=32, negrito=True, fonte=DISPLAY)

    # captura do grafo a esquerda
    img_w = Inches(4.0)
    img_h = Inches(2.5)
    arq = recorte(DADOS["imagem_modulos"], float(img_w) / float(img_h))
    s.shapes.add_picture(arq, MARGEM, Inches(1.75), width=img_w, height=img_h)
    m = caixa(s, MARGEM, Inches(1.75), img_w, img_h, preenche=None, borda=RULE)
    m.line.width = Pt(0.75)
    texto(s, MARGEM, Inches(4.32), img_w, Inches(0.24), DADOS["imagem_modulos_legenda"],
          tam=8, cor=DIM, fonte=MONO)
    texto(s, MARGEM, Inches(4.72), img_w, Inches(1.0),
          "Cada módulo é uma área funcional própria. Juntos, formam a base sobre a qual as "
          "outras aplicações do PLI se apoiam.", tam=11.5, cor=MUTED, espaco=1.2)

    mods = DADOS["modulos_sigma"]
    x0 = MARGEM + img_w + Inches(0.5)
    larg = W - MARGEM - x0
    cols, gap = 3, Inches(0.22)
    cw = (larg - gap * (cols - 1)) / cols
    ch = Inches(1.05)
    for idx, (cod, nome, desc) in enumerate(mods):
        c, r = idx % cols, idx // cols
        x = x0 + c * (cw + gap)
        y = Inches(1.75) + r * (ch + Inches(0.16))
        caixa(s, x, y, cw, ch, preenche=SURFACE, borda=RULE)
        caixa(s, x, y, Pt(3), ch, preenche=BLUE)
        texto(s, x + Inches(0.2), y + Inches(0.11), cw - Inches(0.35), Inches(0.2), cod, tam=8.5, cor=BLUE, fonte=MONO)
        texto(s, x + Inches(0.2), y + Inches(0.32), cw - Inches(0.35), Inches(0.26), nome,
              tam=12.5, negrito=True, fonte=DISPLAY)
        texto(s, x + Inches(0.2), y + Inches(0.58), cw - Inches(0.35), Inches(0.42), desc,
              tam=9.5, cor=MUTED, espaco=1.12)

    rodape(s, DADOS["apps"][0]["url"], "Sistemas Maiores", n, total)
    return s


def slide_integracoes(prs, n, total):
    s = novo_slide(prs)
    faixa(s, GREEN)
    eyebrow(s, MARGEM, Inches(0.5), "Arquitetura", cor=GREEN)
    texto(s, MARGEM, Inches(0.88), Inches(10), Inches(0.6), "Como as peças conversam",
          tam=32, negrito=True, fonte=DISPLAY)
    texto(s, MARGEM, Inches(1.55), Inches(9.8), Inches(0.5),
          "O portfólio não é um conjunto de sistemas isolados: quatro integrações já em produção "
          "ligam sensor, cadastro e decisão.", tam=12.5, cor=MUTED, espaco=1.2)

    bw, bh = Inches(3.0), Inches(1.0)

    def bloco(x, y, cor, nome, sub):
        caixa(s, x, y, bw, bh, preenche=SURFACE, borda=cor)
        caixa(s, x, y, Pt(3.5), bh, preenche=cor)
        texto(s, x + Inches(0.24), y + Inches(0.2), bw - Inches(0.45), Inches(0.3), nome,
              tam=14, negrito=True, fonte=DISPLAY)
        texto(s, x + Inches(0.24), y + Inches(0.54), bw - Inches(0.45), Inches(0.35), sub,
              tam=10, cor=MUTED, espaco=1.1)

    def seta(x1, y1, x2, y2, rotulo):
        con = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
        con.line.color.rgb = DIM
        con.line.width = Pt(1.25)
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        texto(s, Emu(int(mx - Inches(0.62))), Emu(int(my - Inches(0.34))), Inches(1.24), Inches(0.24),
              rotulo.upper(), tam=7.5, cor=DIM, fonte=MONO, alinha=PP_ALIGN.CENTER)

    col1, col2, col3 = MARGEM, MARGEM + Inches(4.3), MARGEM + Inches(8.6)

    texto(s, col1, Inches(2.35), Inches(3.0), Inches(0.22), "SENSORES E COLETA", tam=8, cor=DIM, fonte=MONO)
    bloco(col1, Inches(2.6), BLUE, "PLI-HazardTrack", "chuva MERGE/INPE → risco por trecho")
    bloco(col1, Inches(3.85), BLUE, "PLI Reporta", "relato do cidadão → incidente verificado")

    texto(s, col2, Inches(2.35), Inches(3.0), Inches(0.22), "OPERAÇÃO", tam=8, cor=DIM, fonte=MONO)
    bloco(col2, Inches(3.2), BLUE, "PLI Smart Router", "rota que desvia do risco")

    seta(col1 + bw, Inches(3.1), col2, Inches(3.5), "camadas de risco")
    seta(col1 + bw, Inches(4.35), col2, Inches(3.95), "geojson")

    texto(s, col3, Inches(2.35), Inches(3.0), Inches(0.22), "LEITURA DO TERRITÓRIO", tam=8, cor=DIM, fonte=MONO)
    bloco(col3, Inches(2.6), BLUE, "FAD-Stats 2.0", "IBGE · SECEX · MTE · ANTT")
    bloco(col3, Inches(3.85), AMBER, "Análises Exploratórias", "sinistralidade e densidade da malha")

    texto(s, col1, Inches(5.4), Inches(3.0), Inches(0.22), "PLATAFORMA CENTRAL", tam=8, cor=DIM, fonte=MONO)
    bloco(col1, Inches(5.65), GREEN, "SIGMA-PLI", "metadados, identidade e grafo")

    texto(s, col2, Inches(5.4), Inches(3.0), Inches(0.22), "DECISÃO", tam=8, cor=DIM, fonte=MONO)
    bloco(col2, Inches(5.65), BLUE, "SICARD", "carteira hierarquizada por AHP")

    texto(s, col3, Inches(5.4), Inches(3.0), Inches(0.22), "MÉTODO", tam=8, cor=DIM, fonte=MONO)
    bloco(col3, Inches(5.65), GREEN, "AHP Tool Calculator", "pesos e consistência")

    seta(col1 + bw, Inches(6.15), col2, Inches(6.15), "login e cadastros")
    seta(col3, Inches(6.15), col2 + bw, Inches(6.15), "método AHP")

    rodape(s, "4 integrações em produção", "Arquitetura", n, total)
    return s


def slide_situacao(prs, n, total):
    s = novo_slide(prs)
    faixa(s, BLUE)
    apps = DADOS["apps"]
    no_ar = sum(1 for a in apps if a["status"] == "no ar")
    fora = len(apps) - no_ar

    eyebrow(s, MARGEM, Inches(0.5), "Situação e próximos passos")
    texto(s, MARGEM, Inches(0.88), Inches(10), Inches(0.7), "Onde o portfólio está hoje",
          tam=34, negrito=True, fonte=DISPLAY)

    x = MARGEM
    for valor, rotulo, cor in [(str(no_ar), "no servidor próprio", GREEN),
                               (str(fora), "suspensas no Render", AMBER),
                               ("4", "integrações ativas", TEXT)]:
        texto(s, x, Inches(1.8), Inches(3), Inches(0.7), valor, tam=40, negrito=True, fonte=DISPLAY, cor=cor)
        texto(s, x, Inches(2.42), Inches(3), Inches(0.3), rotulo.upper(), tam=8.5, cor=DIM, fonte=MONO)
        x += Inches(3.4)
    caixa(s, MARGEM, Inches(2.88), LARG, Pt(0.75), preenche=RULE)

    blocos = [
        ("Infraestrutura própria",
         [f"As {no_ar} aplicações dos Sistemas Maiores rodam em containers Docker atrás de um Nginx "
          "na instância EC2 do PLI, cada uma com seu banco PostgreSQL/PostGIS.",
          "É a base madura do portfólio: disponível, versionada e integrada."]),
        ("Hospedagem externa",
         [f"As {fora} aplicações de Análises Exploratórias e Ferramentas estão no Render e hoje "
          "respondem Service Suspended.",
          "Retomá-las exige decidir entre reativar o plano ou migrá-las para a mesma infraestrutura "
          "dos sistemas maiores."]),
        ("Próximo passo",
         ["Concluir a análise aplicação por aplicação e classificar cada uma como produto, módulo, "
          "serviço compartilhado, biblioteca, ferramenta interna ou protótipo.",
          "Só então definir a arquitetura da empresa e a estratégia comercial."]),
    ]
    cw = (LARG - Inches(0.7)) / 3
    x = MARGEM
    for titulo, paragrafos in blocos:
        texto(s, x, Inches(3.3), cw, Inches(0.4), titulo, tam=16, negrito=True, fonte=DISPLAY)
        texto(s, x, Inches(3.85), cw, Inches(2.2), paragrafos, tam=11.5, cor=MUTED, espaco=1.3)
        x += cw + Inches(0.35)

    rodape(s, DADOS["hub_url"], f"{DADOS['org']} · {DADOS['data']}", n, total)
    return s


# ---------------------------------------------------------------- saidas


def gerar_pptx(destino):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    apps = DADOS["apps"]
    total = 4 + (len(apps) - 1) + 2

    n = 1
    slide_capa(prs, n, total); n += 1
    slide_panorama(prs, n, total); n += 1
    slide_app(prs, apps[0], 0, len(apps), n, total); n += 1
    slide_modulos(prs, n, total); n += 1
    for i, app in enumerate(apps[1:], start=1):
        slide_app(prs, app, i, len(apps), n, total); n += 1
    slide_integracoes(prs, n, total); n += 1
    slide_situacao(prs, n, total); n += 1

    prs.save(destino)
    return len(prs.slides._sldIdLst)


def data_uri(caminho):
    dados = (AQUI / caminho).read_bytes()
    return "data:image/jpeg;base64," + base64.b64encode(dados).decode("ascii")


def gerar_html(destino):
    template = (AQUI / "template.html").read_text(encoding="utf-8")
    if "/*DADOS*/" not in template or "/*IMAGENS*/" not in template:
        raise SystemExit("template.html precisa dos marcadores /*DADOS*/ e /*IMAGENS*/")

    imagens = {a["id"]: data_uri(a["imagem"]) for a in DADOS["apps"] if a.get("imagem")}
    imagens["modulos"] = data_uri(DADOS["imagem_modulos"])

    saida = template.replace("/*DADOS*/", json.dumps(DADOS, ensure_ascii=False, indent=2))
    saida = saida.replace("/*IMAGENS*/", json.dumps(imagens))
    destino.write_text(saida, encoding="utf-8")
    return len(imagens)


if __name__ == "__main__":
    html = AQUI / "pli-hub-apresentacao.html"
    pptx = AQUI / "pli-hub-apresentacao.pptx"
    qtd = gerar_html(html)
    n = gerar_pptx(pptx)
    print(f"HTML gerado : {html.name} ({qtd} capturas embutidas, {html.stat().st_size // 1024} KB)")
    print(f"PPTX gerado : {pptx.name} ({n} slides, {pptx.stat().st_size // 1024} KB)")
