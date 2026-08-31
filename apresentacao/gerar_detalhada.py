#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Versao detalhada da apresentacao: so as aplicacoes que tem captura de tela.

Uso:  python gerar_detalhada.py

Saidas (na mesma pasta):
  - pli-hub-detalhada.html   (deck navegavel; imagens embutidas em base64)
  - pli-hub-detalhada.pptx   (PowerPoint 16:9)

Cada aplicacao ocupa dois slides: a vitrine (identidade, resumo e captura
grande) e o detalhamento (como funciona, dados e fontes, integracoes,
numeros e situacao). O conteudo vem de dados.json e de detalhes.json.
"""

import json
import pathlib
import re

import gerar_apresentacao as base
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

AQUI = pathlib.Path(__file__).parent
DADOS = base.DADOS
DETALHES = json.loads((AQUI / "detalhes.json").read_text(encoding="utf-8"))
APPS = [a for a in DADOS["apps"] if a.get("imagem") and a["id"] in DETALHES]

W, H, MARGEM, LARG = base.W, base.H, base.MARGEM, base.LARG
DISPLAY, BODY, MONO = base.DISPLAY, base.BODY, base.MONO
TEXT, MUTED, DIM, RULE = base.TEXT, base.MUTED, base.DIM, base.RULE
SURFACE, GREEN, AMBER, BLUE = base.SURFACE, base.GREEN, base.AMBER, base.BLUE


def cabeca(s, app, i, cor, rotulo=""):
    cat = DADOS["categorias"][app["categoria"]]["nome"]
    base.faixa(s, cor)
    base.eyebrow(s, MARGEM, Inches(0.5), f"Aplicação {i + 1:02d} / {len(APPS)}", cor=cor, w=Inches(2.2))
    texto_cat = cat.upper() + (f" · {rotulo.upper()}" if rotulo else "")
    base.texto(s, MARGEM + Inches(2.3), Inches(0.5), Inches(6.5), Inches(0.24), texto_cat,
               tam=9, cor=DIM, fonte=MONO)
    if app.get("ancoragem"):
        base.texto(s, W - MARGEM - Inches(6.0), Inches(0.5), Inches(4.6), Inches(0.24),
                   app["ancoragem"].upper(), tam=8, cor=cor, fonte=MONO, alinha=PP_ALIGN.RIGHT)
    if app["status"] == "no ar":
        base.pilula(s, W - MARGEM - Inches(0.85), Inches(0.45), "no ar", GREEN, base.DARKINK)
    else:
        base.pilula(s, W - MARGEM - Inches(1.1), Inches(0.45), "suspensa", AMBER, RGBColor(0x2A, 0x14, 0x08))


def slide_vitrine(prs, app, i, n, total):
    """Vitrine: identidade e stack a esquerda; captura grande e entregas a direita."""
    s = base.novo_slide(prs)
    cor = base.cor_cat(app["categoria"])
    cabeca(s, app, i, cor)

    esq_w = Inches(4.55)
    base.texto(s, MARGEM, Inches(1.15), esq_w, Inches(0.8), app["nome"],
               tam=31, negrito=True, fonte=DISPLAY, espaco=0.95)
    base.texto(s, MARGEM, Inches(1.88), esq_w, Inches(0.5), app["extenso"], tam=11, cor=MUTED, espaco=1.15)
    base.texto(s, MARGEM, Inches(2.62), esq_w, Inches(1.8), app["resumo"], tam=12.5, espaco=1.32)

    base.caixa(s, MARGEM, Inches(4.72), Pt(3), Inches(0.38), preenche=cor)
    base.texto(s, MARGEM + Inches(0.15), Inches(4.76), esq_w, Inches(0.34), app["papel"],
               tam=13, negrito=True, fonte=DISPLAY, cor=cor)

    base.caixa(s, MARGEM, Inches(5.42), esq_w, Pt(0.75), preenche=RULE)
    base.chips(s, MARGEM, Inches(5.62), esq_w, app["stack"])

    # captura a direita
    x = MARGEM + esq_w + Inches(0.45)
    w = W - MARGEM - x
    h = Inches(3.1)
    y = Inches(1.42)
    arq = base.recorte(app["imagem"], float(w) / float(h))
    barra = base.caixa(s, x, y - Inches(0.26), w, Inches(0.26), preenche=base.SURFACE2, borda=RULE)
    tf = barra.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = "   " + app["url"]
    r.font.size = Pt(8)
    r.font.name = MONO
    r.font.color.rgb = DIM
    for k in range(3):
        d = base.caixa(s, x + Inches(0.1) + k * Inches(0.14), y - Inches(0.2),
                       Inches(0.08), Inches(0.08), preenche=RULE)
        d.line.fill.background()
    s.shapes.add_picture(arq, x, y, width=w, height=h)
    moldura = base.caixa(s, x, y, w, h, preenche=None, borda=RULE)
    moldura.line.width = Pt(0.75)
    base.texto(s, x, y + h + Inches(0.08), w, Inches(0.24), app["imagem_legenda"], tam=8, cor=DIM, fonte=MONO)

    base.texto(s, x, Inches(5.0), w, Inches(0.22), "O QUE ELA ENTREGA", tam=8.5, cor=DIM, fonte=MONO)
    base.lista_destaques(s, x, Inches(5.3), w, app["destaques"], cor, tam=10, chars=82)

    base.rodape(s, app["url"], DADOS["categorias"][app["categoria"]]["nome"], n, total)
    return s


def slide_detalhe(prs, app, i, n, total):
    s = base.novo_slide(prs)
    cor = base.cor_cat(app["categoria"])
    d = DETALHES[app["id"]]
    cabeca(s, app, i, cor, "detalhamento")

    base.texto(s, MARGEM, Inches(0.98), Inches(8), Inches(0.5), f"{app['nome']} · por dentro",
               tam=24, negrito=True, fonte=DISPLAY)

    # chamada de largura total: o papel da aplicacao dentro do plano
    barra_papel = base.caixa(s, MARGEM, Inches(1.52), Pt(3), Inches(0.5), preenche=cor)
    barra_papel.line.fill.background()
    base.texto(s, MARGEM + Inches(0.18), Inches(1.52), LARG - Inches(0.3), Inches(0.5),
               d["no_plano"], tam=10.5, cor=MUTED, espaco=1.2)

    gap = Inches(0.45)
    c1 = Inches(4.9)
    c2 = Inches(4.1)
    c3 = LARG - c1 - c2 - 2 * gap
    x1 = MARGEM
    x2 = x1 + c1 + gap
    x3 = x2 + c2 + gap
    y0 = Inches(2.28)

    def titulo(x, y, w, txt):
        base.texto(s, x, y, w, Inches(0.2), txt.upper(), tam=8, cor=DIM, fonte=MONO)
        linha = base.caixa(s, x, y + Inches(0.24), w, Pt(0.75), preenche=RULE)
        linha.line.fill.background()
        return y + Inches(0.42)

    # coluna 1 — como funciona
    y = titulo(x1, y0, c1, "Como funciona")
    base.lista_destaques(s, x1, y, c1, d["como_funciona"], cor, tam=10, chars=60)

    # coluna 2 — dados e integracoes
    y = titulo(x2, y0, c2, "Dados e fontes")
    y = base.lista_destaques(s, x2, y, c2, d["dados"], cor, tam=10, chars=50)
    y = titulo(x2, y + Inches(0.25), c2, "Integrações e publicação")
    base.lista_destaques(s, x2, y, c2, d["integracoes"], cor, tam=10, chars=50)

    # coluna 3 — numeros e situacao
    y = titulo(x3, y0, c3, "Números")
    for valor, rotulo in d["numeros"]:
        base.texto(s, x3, y, c3, Inches(0.28), valor, tam=15, negrito=True, fonte=DISPLAY, cor=cor)
        base.texto(s, x3, y + Inches(0.25), c3, Inches(0.2), rotulo.upper(), tam=7.5, cor=DIM, fonte=MONO)
        linha = base.caixa(s, x3, y + Inches(0.48), c3, Pt(0.75), preenche=RULE)
        linha.line.fill.background()
        y += Inches(0.62)

    y = titulo(x3, y + Inches(0.14), c3, "Situação")
    base.texto(s, x3, y, c3, Inches(1.3), d["situacao"], tam=9, cor=MUTED, espaco=1.24)

    base.rodape(s, app["url"], "  ·  ".join(app["stack"][:4]), n, total)
    return s


def slide_capa(prs, n, total):
    s = base.novo_slide(prs)
    base.faixa(s, BLUE)
    base.caixa(s, W - Inches(5.4), 0, Inches(5.4), H, preenche=RGBColor(0x08, 0x22, 0x30))
    base.caixa(s, W - Inches(5.4), 0, Pt(1), H, preenche=RULE)

    base.eyebrow(s, MARGEM, Inches(0.95), f"{DADOS['org']} · {DADOS['data']}")
    base.texto(s, MARGEM, Inches(1.5), Inches(7.2), Inches(1.2), "PLI Hub",
               tam=76, negrito=True, fonte=DISPLAY, espaco=0.85)
    base.texto(s, MARGEM, Inches(2.62), Inches(7.2), Inches(1.0), "em detalhe",
               tam=76, negrito=True, fonte=DISPLAY, espaco=0.85, cor=BLUE)
    base.texto(s, MARGEM, Inches(4.0), Inches(6.2), Inches(1.0),
               f"As {len(APPS)} aplicações com tela registrada, uma a uma: o que fazem, "
               "como funcionam e de que dados vivem.", tam=17, cor=MUTED, espaco=1.25)

    y = Inches(5.35)
    for rotulo, valor in [
        ("Recorte", f"{len(APPS)} de {len(DADOS['apps'])} aplicações do hub"),
        ("Central de acesso", DADOS["hub_url"]),
        ("Responsável técnico", DADOS["autor"]),
    ]:
        base.texto(s, MARGEM, y, Inches(6.5), Inches(0.22), rotulo.upper(), tam=8.5, cor=DIM, fonte=MONO)
        base.texto(s, MARGEM, y + Inches(0.22), Inches(6.5), Inches(0.3), valor, tam=12.5)
        y += Inches(0.6)

    x = W - Inches(5.4) + Inches(0.7)
    cw = Inches(4.0)
    yy = Inches(1.5)
    base.texto(s, x, yy, cw, Inches(0.24), "NESTE DOCUMENTO", tam=9, cor=DIM, fonte=MONO)
    yy += Inches(0.45)
    for cid, cat in DADOS["categorias"].items():
        lista = [a for a in APPS if a["categoria"] == cid]
        if not lista:
            continue
        base.caixa(s, x, yy, cw, Pt(2.5), preenche=base.hex_rgb(cat["cor"]))
        base.texto(s, x, yy + Inches(0.12), cw, Inches(0.28), cat["nome"],
                   tam=13, negrito=True, fonte=DISPLAY, cor=base.hex_rgb(cat["cor"]))
        yy += Inches(0.42)
        for a in lista:
            base.texto(s, x, yy, cw, Inches(0.24), a["nome"], tam=11, cor=MUTED)
            yy += Inches(0.27)
        yy += Inches(0.22)
    return s


def slide_recorte(prs, n, total):
    s = base.novo_slide(prs)
    base.faixa(s, BLUE)
    no_ar = sum(1 for a in APPS if a["status"] == "no ar")
    base.eyebrow(s, MARGEM, Inches(0.5), "Recorte deste documento")
    base.texto(s, MARGEM, Inches(0.88), Inches(9), Inches(0.7), "Só o que já tem tela",
               tam=34, negrito=True, fonte=DISPLAY)
    base.texto(s, MARGEM, Inches(1.65), Inches(11.5), Inches(0.8),
               f"Das {len(DADOS['apps'])} aplicações do PLI Hub, {len(APPS)} têm registro visual da interface — "
               "capturadas ao vivo no servidor de produção ou recuperadas das capturas do relatório de "
               "atividades. São essas que este documento detalha, com dois slides cada: a vitrine e o "
               "funcionamento por dentro.", tam=12.5, cor=MUTED, espaco=1.25)

    x = MARGEM
    for valor, rotulo, cor in [(str(len(APPS)), "aplicações detalhadas", TEXT),
                               (str(no_ar), "no ar", GREEN),
                               (str(len(APPS) - no_ar), "suspensas", AMBER),
                               (str(len(APPS) * 2), "slides de conteúdo", TEXT)]:
        base.texto(s, x, Inches(2.75), Inches(2.8), Inches(0.7), valor, tam=38, negrito=True, fonte=DISPLAY, cor=cor)
        base.texto(s, x, Inches(3.35), Inches(2.8), Inches(0.3), rotulo.upper(), tam=8.5, cor=DIM, fonte=MONO)
        x += Inches(2.9)
    base.caixa(s, MARGEM, Inches(3.8), LARG, Pt(0.75), preenche=RULE)

    col_w = (LARG - Inches(0.7)) / 3
    x = MARGEM
    for cid, cat in DADOS["categorias"].items():
        lista = [a for a in APPS if a["categoria"] == cid]
        if not lista:
            continue
        cor = base.hex_rgb(cat["cor"])
        base.caixa(s, x, Inches(4.15), col_w, Pt(3), preenche=cor)
        base.texto(s, x, Inches(4.35), col_w, Inches(0.5), cat["nome"],
                   tam=15, negrito=True, fonte=DISPLAY, cor=cor, espaco=1.05)
        y = Inches(5.0)
        for k, a in enumerate(lista, 1):
            base.texto(s, x, y, Inches(0.3), Inches(0.26), f"{k:02d}", tam=9, cor=DIM, fonte=MONO)
            base.texto(s, x + Inches(0.42), y - Inches(0.02), col_w - Inches(0.8), Inches(0.26), a["nome"], tam=12)
            ponto = base.caixa(s, x + col_w - Inches(0.22), y + Inches(0.05), Inches(0.1), Inches(0.1),
                               preenche=GREEN if a["status"] == "no ar" else AMBER)
            ponto.line.fill.background()
            base.caixa(s, x, y + Inches(0.3), col_w - Inches(0.1), Pt(0.75), preenche=RULE)
            y += Inches(0.42)
        x += col_w + Inches(0.35)

    base.rodape(s, DADOS["hub_url"], "Recorte", n, total)
    return s


# ---------------------------------------------------------------- saidas


def gerar_pptx(destino):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    total = 3 + len(APPS) * 2 + 3

    n = 1
    slide_capa(prs, n, total); n += 1
    base.slide_contexto(prs, n, total); n += 1
    slide_recorte(prs, n, total); n += 1
    for i, app in enumerate(APPS):
        slide_vitrine(prs, app, i, n, total); n += 1
        slide_detalhe(prs, app, i, n, total); n += 1
    base.slide_cadeia(prs, n, total); n += 1
    base.slide_integracoes(prs, n, total); n += 1
    base.slide_fontes(prs, n, total)

    prs.save(destino)
    return len(prs.slides._sldIdLst)


def gerar_html(destino):
    completo = (AQUI / "template.html").read_text(encoding="utf-8")
    estilo_base = re.search(r"<style>(.*?)</style>", completo, re.S).group(1)

    template = (AQUI / "template_detalhe.html").read_text(encoding="utf-8")
    for marcador in ("/*ESTILO_BASE*/", "/*DADOS*/", "/*DETALHES*/", "/*IMAGENS*/", "/*COMUM*/"):
        if marcador not in template:
            raise SystemExit(f"template_detalhe.html precisa do marcador {marcador}")

    imagens = {a["id"]: base.data_uri(a["imagem"]) for a in APPS}

    comum = (AQUI / "comum.js").read_text(encoding="utf-8")
    saida = template.replace("/*ESTILO_BASE*/", estilo_base)
    saida = saida.replace("/*COMUM*/", comum)
    saida = saida.replace("/*DADOS*/", json.dumps(DADOS, ensure_ascii=False, indent=2))
    saida = saida.replace("/*DETALHES*/", json.dumps(DETALHES, ensure_ascii=False, indent=2))
    saida = saida.replace("/*IMAGENS*/", json.dumps(imagens))
    destino.write_text(saida, encoding="utf-8")
    return len(imagens)


if __name__ == "__main__":
    html = AQUI / "pli-hub-detalhada.html"
    pptx = AQUI / "pli-hub-detalhada.pptx"
    qtd = gerar_html(html)
    n = gerar_pptx(pptx)
    print(f"Aplicacoes detalhadas: {len(APPS)} -> {', '.join(a['nome'] for a in APPS)}")
    print(f"HTML gerado : {html.name} ({qtd} capturas, {html.stat().st_size // 1024} KB)")
    print(f"PPTX gerado : {pptx.name} ({n} slides, {pptx.stat().st_size // 1024} KB)")
